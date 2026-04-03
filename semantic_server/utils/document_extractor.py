import fitz
import docx
from pptx import Presentation


def extract_text_from_pdf(path):

    doc = fitz.open(path)

    text = ""

    for page in doc:
        text += page.get_text()

    return text


def extract_text_from_docx(path):

    doc = docx.Document(path)

    text = "\n".join([p.text for p in doc.paragraphs])

    return text


def extract_text_from_pptx(path):

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text